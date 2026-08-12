import importlib
import importlib.util
import re
import subprocess
import time
import unittest
from pathlib import Path
from tempfile import TemporaryDirectory
from unittest.mock import patch
from urllib.parse import unquote, urljoin, urlparse

from openpyxl import load_workbook
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

from tools import build_research_model_v03_workbook as workbook_builder


build_workbook = workbook_builder.build_workbook


ROOT = Path(__file__).resolve().parents[1]
SHEET_ORDER = [
    "README",
    "구성개념",
    "명제추적",
    "논문별기여",
    "사건코딩",
    "관계망",
    "부정사례",
    "윤리체크",
]
EVENT_COLUMNS = [
    "event_code",
    "actor_role",
    "trigger",
    "threat",
    "role_complementarity",
    "activated_connectivity",
    "voice_safety",
    "voice_efficacy",
    "response_codes",
    "dominant_path",
    "substantive_handling",
    "problem_resolution",
    "rival_explanation",
    "negative_case",
    "dv1_decision",
    "dv1_primary_source",
    "dv1_corroboration",
    "dv2_decision",
    "dv2_primary_source",
    "dv2_corroboration",
    "dv3_decision",
    "dv3_primary_source",
    "dv3_corroboration",
    "dv4_decision",
    "dv4_primary_source",
    "dv4_corroboration",
    "divergence_vector",
    "divergence_event_decision",
    "divergence_aggregation",
    "divergence_uncertainty",
    "p6_use_decision",
    "p7_use_decision",
]
ETHICS_IDS = [f"ETH{number:02d}" for number in range(1, 15)]
SLIDE_SEMANTIC_CONTRACT = (
    {
        "title": "분산된 AI 변화주도자 네트워크",
        "facts": ("공식 리더와 비공식 AI 챔피언", "P1-P7은 검증 전 명제"),
        "sources": ("site/downloads/research-design/research-model-v0.3.md",),
    },
    {
        "title": "공식 도입과 비공식 실험이 같은 사건에 겹친다",
        "facts": ("Shadow AI", "구체적 AI 사건"),
        "sources": (
            "https://doi.org/10.1016/j.infsof.2025.107805",
            "https://doi.org/10.1007/978-3-032-22375-3_18",
        ),
    },
    {
        "title": "다섯 논문은 현상·과정·경계의 서로 다른 근거다",
        "facts": (
            "Kemell et al. (2025)",
            "다섯 논문을 합쳐도 P1-P7의 직접 검증 근거가 되지 않는다",
        ),
        "sources": (
            "https://doi.org/10.1016/j.hrmr.2024.101075",
            "https://doi.org/10.5465/amj.2009.0891",
            "https://doi.org/10.1287/mnsc.1120.1583",
        ),
    },
    {
        "title": "AI 사건을 중심으로 세 분석 수준을 분리한다",
        "facts": ("팀 맥락", "단위 계약: 전달 성공은 연결성"),
        "sources": ("site/downloads/research-design/pilot-protocol-and-codingbook-v0.3.md",),
    },
    {
        "title": "역할 보완성은 권한과 현장 기능의 실제 분담이다",
        "facts": ("실행된 역할 보완성", "직급·친분·개인 열의 또는 해결 성과"),
        "sources": ("https://doi.org/10.1016/j.hrmr.2024.101075",),
    },
    {
        "title": "네트워크 연결성은 네 관계의 도달·확인을 묻는다",
        "facts": ("수신·확인 근거", "발신자 보고만으로 수신을 확정하지 않는다"),
        "sources": ("tools/build_research_model_v03_workbook.py",),
    },
    {
        "title": "안전성과 효능감은 서로 다른 발언 판단이다",
        "facts": ("발언 안전성", "같은 구성개념이 아니다"),
        "sources": (
            "site/downloads/research-design/research-model-v0.3.png",
            "site/downloads/research-design/construct-dictionary-v0.3.md",
        ),
    },
    {
        "title": "한 사건의 반응은 복수 행동과 지배 경로로 코딩한다",
        "facts": ("건설적 표면화", "우회와 Shadow AI는"),
        "sources": ("site/downloads/research-design/pilot-protocol-and-codingbook-v0.3.md",),
    },
    {
        "title": "변화 발산성이 유리한 관계 구조를 바꿀 수 있다",
        "facts": ("내부 응집·강한 관계", "P7 경계:"),
        "sources": (
            "https://doi.org/10.5465/amj.2009.0891",
            "https://doi.org/10.1287/mnsc.1120.1583",
        ),
    },
    {
        "title": "P1-P7은 경쟁 설명과 함께 검증할 명제다",
        "facts": ("P1", "판정 규칙:"),
        "sources": ("site/downloads/research-design/proposition-traceability-v0.3.md",),
    },
    {
        "title": "탐색적 순차 혼합연구가 명제를 단계적으로 거른다",
        "facts": ("프레임워크 개발", "특정 표본 수나 SEM 유형을 고정하지 않는다"),
        "sources": (
            "site/downloads/research-design/research-model-v0.3.md",
            "site/downloads/research-design/pilot-protocol-and-codingbook-v0.3.md",
        ),
    },
    {
        "title": "현장 진입 전 윤리 게이트와 다음 결정을 확인한다",
        "facts": ("IRB + 기업 보안·법무 승인", "NotebookLM 산출물은 근거가 아니다"),
        "sources": (
            "docs/reviews/2026-08-12-research-model-v0.3-review.md",
            "site/downloads/research-design/pilot-protocol-and-codingbook-v0.3.md",
        ),
    },
)
RELATION_BOOLEAN_COLUMNS = [
    "advice_relation",
    "trust_relation",
    "approval_exception_relation",
    "risk_escalation_relation",
]


def normalized_pdf_text(page) -> str:
    """Recover Korean text from LibreOffice PDFs whose font map exposes CP949 bytes."""

    decoded = []
    byte_chars = []

    def flush_byte_chars() -> None:
        if not byte_chars:
            return
        raw = "".join(byte_chars).encode("latin1")
        try:
            decoded.append(raw.decode("cp949"))
        except UnicodeDecodeError:
            decoded.append(raw.decode("latin1"))
        byte_chars.clear()

    for character in page.extract_text() or "":
        if ord(character) <= 255:
            byte_chars.append(character)
        else:
            flush_byte_chars()
            decoded.append(character)
    flush_byte_chars()
    return " ".join("".join(decoded).split())
RELATION_TYPES = ["advice", "trust", "approval_exception", "risk_escalation"]
RELATION_EVIDENCE_SUFFIXES = [
    "sender_report_delivery",
    "receipt_confirmation_evidence",
    "recipient_stance",
]


def headers(worksheet):
    return [cell.value for cell in worksheet[1]]


def workbook_snapshot(workbook):
    return tuple(
        (
            worksheet.title,
            worksheet.freeze_panes,
            worksheet.auto_filter.ref,
            tuple(tuple(cell.value for cell in row) for row in worksheet.iter_rows()),
            tuple(
                (cell.coordinate, cell.hyperlink.target)
                for row in worksheet.iter_rows()
                for cell in row
                if cell.hyperlink
            ),
            tuple(
                (validation.type, validation.formula1, str(validation.sqref))
                for validation in worksheet.data_validations.dataValidation
            ),
        )
        for worksheet in workbook.worksheets
    )


class ResearchModelV03AssetTests(unittest.TestCase):
    def build(self, directory, name="workbook.xlsx"):
        path = Path(directory) / name
        build_workbook(path)
        return path, load_workbook(path, read_only=False, data_only=False)

    def test_v03_workbook_structure(self):
        with TemporaryDirectory() as tmp:
            _, workbook = self.build(tmp)

            self.assertEqual(workbook.sheetnames, SHEET_ORDER)
            self.assertEqual(workbook["구성개념"].max_row, 13)
            self.assertEqual(workbook["명제추적"].max_row, 8)
            self.assertEqual(workbook["논문별기여"].max_row, 6)
            self.assertEqual(
                [workbook["명제추적"].cell(row=row, column=1).value for row in range(2, 9)],
                [f"P{number}" for number in range(1, 8)],
            )
            self.assertEqual(
                {workbook["명제추적"].cell(row=row, column=3).value for row in range(2, 9)},
                {"검증 전"},
            )
            self.assertIn("실행된 역할 보완성", {cell.value for cell in workbook["구성개념"]["A"]})
            self.assertIn("Kemell et al. (2025)", {cell.value for cell in workbook["논문별기여"]["A"]})

    def test_blank_templates_preserve_event_and_network_boundaries(self):
        with TemporaryDirectory() as tmp:
            _, workbook = self.build(tmp)

            self.assertEqual(headers(workbook["사건코딩"]), EVENT_COLUMNS)
            self.assertEqual(workbook["사건코딩"].max_row, 1)

            relation_headers = headers(workbook["관계망"])
            self.assertEqual(relation_headers[:4], ["event_code", "respondent_code", "alter_code", "alter_role_category"])
            self.assertEqual(
                [column for column in relation_headers if column in RELATION_BOOLEAN_COLUMNS],
                RELATION_BOOLEAN_COLUMNS,
            )
            self.assertEqual(workbook["관계망"].max_row, 1)
            self.assertNotIn("real_name", relation_headers)
            self.assertFalse(any("실명" in str(value) or "name" == str(value).lower() for value in relation_headers))

            self.assertEqual(workbook["부정사례"].max_row, 1)
            self.assertTrue(
                all(
                    cell.value in (None, "")
                    for row in workbook["윤리체크"].iter_rows(min_row=2)
                    for cell in row[2:5]
                )
            )

    def test_event_sheet_binds_divergence_dimensions_and_p6_p7_decisions(self):
        with TemporaryDirectory() as tmp:
            _, workbook = self.build(tmp)
            worksheet = workbook["사건코딩"]

            self.assertEqual(headers(worksheet), EVENT_COLUMNS)
            validations = {
                validation.formula1: {str(cell_range) for cell_range in validation.ranges.ranges}
                for validation in worksheet.data_validations.dataValidation
            }
            self.assertEqual(
                validations,
                {
                    '"유지,국소 조정,경계 재구성,자료 부족,자료 충돌"': {
                        "O2:O500",
                        "R2:R500",
                        "U2:U500",
                        "X2:X500",
                    },
                    '"고발산 후보,저발산 후보,혼합,판정 유보"': {"AB2:AB500"},
                    '"차원 벡터 보존 (합산·평균 금지)"': {"AC2:AC500"},
                    '"경계 연결 평가,내부 응집·강한 관계 평가,차원별 분기,판정 유보"': {"AE2:AE500"},
                    '"고발산 입력 평가,입력 불충족,고발산 입력 별도 확인,판정 유보"': {"AF2:AF500"},
                },
            )
            for validation in worksheet.data_validations.dataValidation:
                self.assertTrue(validation.showErrorMessage)
                self.assertEqual(validation.errorStyle, "stop")

    def test_ethics_safeguards_keep_canonical_ids_when_prose_becomes_bullets(self):
        pilot_path = ROOT / "site/downloads/research-design/pilot-protocol-and-codingbook-v0.3.md"
        original = pilot_path.read_text(encoding="utf-8")
        safeguards = [row[1] for row in workbook_builder._ethics_rows()]
        self.assertEqual(len(safeguards), 14)

        mutated = original
        for safeguard in safeguards[:9]:
            before = f"\n{safeguard}\n"
            after = f"\n- {safeguard}\n"
            self.assertIn(before, mutated)
            mutated = mutated.replace(before, after, 1)

        with patch.object(workbook_builder, "_read", return_value=mutated):
            rows = workbook_builder._ethics_rows()

        self.assertEqual([row[0] for row in rows], ETHICS_IDS)
        self.assertEqual([row[1] for row in rows], safeguards)

        with TemporaryDirectory() as tmp:
            _, workbook = self.build(tmp)
            worksheet = workbook["윤리체크"]
            self.assertEqual(worksheet.max_row, 15)
            self.assertEqual([cell.value for cell in worksheet["A"][1:]], ETHICS_IDS)

    def test_relation_sheet_preserves_delivery_receipt_and_stance_per_relation(self):
        with TemporaryDirectory() as tmp:
            path, workbook = self.build(tmp)
            worksheet = workbook["관계망"]
            relation_headers = headers(worksheet)

            expected_evidence_columns = [
                f"{relation}_{suffix}"
                for relation in RELATION_TYPES
                for suffix in RELATION_EVIDENCE_SUFFIXES
            ]
            self.assertTrue(set(expected_evidence_columns).issubset(relation_headers))
            self.assertFalse(set(RELATION_EVIDENCE_SUFFIXES).intersection(relation_headers))
            self.assertEqual(len(set(expected_evidence_columns)), 12)

            values = {
                "event_code": "EV-TEST",
                "respondent_code": "R-TEST",
                "alter_code": "A-TEST",
                "advice_relation": True,
                "trust_relation": True,
                "advice_sender_report_delivery": "조언 전달 보고",
                "advice_receipt_confirmation_evidence": "조언 수신 확인",
                "advice_recipient_stance": "수용",
                "trust_sender_report_delivery": "신뢰 전달 보고",
                "trust_receipt_confirmation_evidence": "신뢰 수신 미확인",
                "trust_recipient_stance": "보류",
            }
            for column_name, value in values.items():
                worksheet.cell(row=2, column=relation_headers.index(column_name) + 1).value = value
            workbook.save(path)

            rebuilt = load_workbook(path, read_only=False, data_only=False)["관계망"]
            saved = dict(zip(headers(rebuilt), [cell.value for cell in rebuilt[2]]))
            for column_name, value in values.items():
                self.assertEqual(saved[column_name], value)
            self.assertNotEqual(
                saved["advice_receipt_confirmation_evidence"],
                saved["trust_receipt_confirmation_evidence"],
            )

    def test_headers_filters_freeze_panes_and_layout_are_usable(self):
        with TemporaryDirectory() as tmp:
            _, workbook = self.build(tmp)

            for worksheet in workbook.worksheets:
                self.assertEqual(worksheet.freeze_panes, "A2", worksheet.title)
                self.assertEqual(
                    worksheet.auto_filter.ref,
                    f"A1:{worksheet.cell(row=worksheet.max_row, column=worksheet.max_column).coordinate}",
                    worksheet.title,
                )
                self.assertFalse(worksheet.sheet_view.showGridLines, worksheet.title)
                for cell in worksheet[1]:
                    self.assertEqual(cell.fill.fgColor.rgb, "FF17365D", f"{worksheet.title}!{cell.coordinate}")
                    self.assertEqual(cell.font.name, "Malgun Gothic", f"{worksheet.title}!{cell.coordinate}")
                    self.assertTrue(cell.alignment.wrap_text, f"{worksheet.title}!{cell.coordinate}")
                    self.assertEqual(cell.border.right.style, "thin", f"{worksheet.title}!{cell.coordinate}")
                    width = worksheet.column_dimensions[cell.column_letter].width
                    self.assertGreaterEqual(width, 12, f"{worksheet.title}!{cell.column_letter}")
                    self.assertLessEqual(width, 60, f"{worksheet.title}!{cell.column_letter}")
            self.assertEqual(workbook["README"]["A2"].fill.fgColor.rgb, "FFEAF2F8")
            self.assertGreater(
                workbook["논문별기여"].row_dimensions[3].height,
                180,
                "긴 직접 근거가 인쇄 렌더에서 잘리지 않아야 한다.",
            )

    def test_dense_research_sheets_use_readable_multi_page_print_layouts(self):
        with TemporaryDirectory() as tmp:
            _, workbook = self.build(tmp)

            title_columns = {
                "구성개념": "$A:$A",
                "명제추적": "$A:$A",
                "논문별기여": "$A:$A",
                "사건코딩": "$A:$A",
                "관계망": "$A:$D",
            }
            for sheet_name, expected_title_columns in title_columns.items():
                worksheet = workbook[sheet_name]
                self.assertFalse(worksheet.sheet_properties.pageSetUpPr.fitToPage, sheet_name)
                self.assertIsNone(worksheet.page_setup.fitToWidth, sheet_name)
                self.assertIsNone(worksheet.page_setup.fitToHeight, sheet_name)
                self.assertEqual(worksheet.page_setup.scale, 100, sheet_name)
                self.assertEqual(str(worksheet.page_setup.paperSize), str(worksheet.PAPERSIZE_A3), sheet_name)
                self.assertEqual(worksheet.page_setup.orientation, worksheet.ORIENTATION_LANDSCAPE, sheet_name)
                self.assertEqual(worksheet.print_title_rows, "$1:$1", sheet_name)
                self.assertEqual(worksheet.print_title_cols, expected_title_columns, sheet_name)

    def test_validations_and_paper_lab_links_are_operational(self):
        with TemporaryDirectory() as tmp:
            path, workbook = self.build(tmp)

            proposition_validations = workbook["명제추적"].data_validations.dataValidation
            self.assertEqual(len(proposition_validations), 2)
            self.assertEqual(
                {validation.formula1 for validation in proposition_validations},
                {'"검증 전,수정,분기,기각"', '"미검토,유지,수정,분기,기각"'},
            )
            self.assertEqual(
                {validation.formula1 for validation in workbook["관계망"].data_validations.dataValidation},
                {'"TRUE,FALSE"'},
            )
            self.assertEqual(
                {validation.formula1 for validation in workbook["부정사례"].data_validations.dataValidation},
                {'"미검토,유지,수정,분기,기각"'},
            )

            links = [
                cell.hyperlink.target
                for sheet_name in ("README", "구성개념", "명제추적", "논문별기여", "윤리체크")
                for row in workbook[sheet_name].iter_rows()
                for cell in row
                if cell.hyperlink
            ]
            self.assertGreaterEqual(len(links), 29)
            for target in links:
                parsed = urlparse(target)
                self.assertEqual((parsed.scheme, parsed.netloc), ("https", "beaten-to-it.github.io"), target)
                self.assertTrue(parsed.path.startswith("/paper/downloads/"), target)
                self.assertEqual(urljoin(path.as_uri(), target), target, target)
                source = ROOT / "site" / unquote(parsed.path.removeprefix("/paper/"))
                self.assertTrue(source.is_file(), target)

    def test_boolean_relation_validation_uses_stop_error_alerts(self):
        with TemporaryDirectory() as tmp:
            _, workbook = self.build(tmp)
            validations = workbook["관계망"].data_validations.dataValidation

            self.assertEqual(len(validations), 1)
            validation = validations[0]
            self.assertEqual(validation.formula1, '"TRUE,FALSE"')
            self.assertEqual(
                {str(cell_range) for cell_range in validation.ranges.ranges},
                {"E2:E500", "F2:F500", "G2:G500", "H2:H500"},
            )
            self.assertTrue(validation.showErrorMessage)
            self.assertEqual(validation.errorStyle, "stop")
            self.assertEqual(validation.errorTitle, "허용되지 않은 관계 값")
            self.assertIn("TRUE 또는 FALSE", validation.error)

    def test_rebuild_has_stable_structural_values(self):
        with TemporaryDirectory() as tmp:
            _, first = self.build(tmp, "first.xlsx")
            _, second = self.build(tmp, "second.xlsx")

            self.assertEqual(workbook_snapshot(first), workbook_snapshot(second))

    def test_delayed_rebuild_is_byte_deterministic(self):
        with TemporaryDirectory() as tmp:
            first_path = Path(tmp) / "first.xlsx"
            second_path = Path(tmp) / "second.xlsx"
            build_workbook(first_path)
            time.sleep(2.1)
            build_workbook(second_path)

            self.assertEqual(first_path.read_bytes(), second_path.read_bytes())

    def test_v03_visual_package_is_editable_sourced_and_claim_safe(self):
        module_name = "tools.build_research_model_v03_media"
        self.assertIsNotNone(
            importlib.util.find_spec(module_name),
            "The v0.3 media builder is not implemented.",
        )
        build_media = importlib.import_module(module_name).build_media

        with TemporaryDirectory() as tmp:
            diagram_path = Path(tmp) / "model.png"
            deck_path = Path(tmp) / "deck.pptx"
            build_media(diagram_path, deck_path)
            second_diagram_path = Path(tmp) / "model-second.png"
            second_deck_path = Path(tmp) / "deck-second.pptx"
            build_media(second_diagram_path, second_deck_path)

            with Image.open(diagram_path) as diagram:
                self.assertEqual(diagram.size, (2400, 1600))
                nonwhite = ImageChops.difference(diagram.convert("RGB"), Image.new("RGB", diagram.size, "white"))
                content_bounds = nonwhite.getbbox()
                self.assertIsNotNone(content_bounds)
                left, top, right, bottom = content_bounds
                self.assertGreaterEqual(min(left, top, 2400 - right, 1600 - bottom), 20)
            self.assertEqual(diagram_path.read_bytes(), second_diagram_path.read_bytes())
            self.assertEqual(deck_path.read_bytes(), second_deck_path.read_bytes())

            deck = Presentation(deck_path)
            self.assertEqual(len(deck.slides), 12)
            slide_texts = [
                " ".join(
                    "\n".join(
                        shape.text
                        for shape in slide.shapes
                        if getattr(shape, "has_text_frame", False)
                    ).split()
                )
                for slide in deck.slides
            ]
            slide_titles = []
            for index, slide in enumerate(deck.slides, start=1):
                title_shapes = [shape for shape in slide.shapes if shape.name == "Slide Title"]
                self.assertEqual(len(title_shapes), 1, f"slide {index} title identity")
                slide_titles.append(title_shapes[0].text)
            expected_titles = [contract["title"] for contract in SLIDE_SEMANTIC_CONTRACT]
            self.assertEqual(slide_titles, expected_titles)
            self.assertEqual(len(set(slide_titles)), len(SLIDE_SEMANTIC_CONTRACT))

            for index, (slide, slide_text, contract) in enumerate(
                zip(deck.slides, slide_texts, SLIDE_SEMANTIC_CONTRACT, strict=True),
                start=1,
            ):
                for fact in contract["facts"]:
                    self.assertIn(fact, slide_text, f"slide {index} required fact")
                source_notes = slide.notes_slide.notes_text_frame.text
                for source in contract["sources"]:
                    self.assertIn(source, source_notes, f"slide {index} required source")
            all_text = "\n".join(slide_texts)

            for phrase in (
                "역할 보완성",
                "네트워크 연결성",
                "발언 안전성",
                "발언 효능감",
                "변화 발산성",
                "검증 전 명제",
            ):
                self.assertIn(phrase, all_text)
            self.assertIn("NotebookLM 산출물은 근거가 아니다", all_text)
            self.assertIn("현장조사는 시작하지 않았다", all_text)
            self.assertIsNone(re.search(r"검증됐다|입증됐다", all_text))

            for index, slide in enumerate(deck.slides, start=1):
                source_notes = slide.notes_slide.notes_text_frame.text
                self.assertRegex(source_notes, r"(?s)\[Sources\]\s*\n\S+", f"slide {index}")
                for shape in slide.shapes:
                    self.assertGreaterEqual(shape.left, 0, f"slide {index}: {shape.name}")
                    self.assertGreaterEqual(shape.top, 0, f"slide {index}: {shape.name}")
                    self.assertGreaterEqual(
                        shape.left,
                        457200,
                        f"slide {index}: {shape.name} violates the 0.5-inch left safe area",
                    )
                    self.assertLessEqual(
                        shape.left + shape.width,
                        deck.slide_width,
                        f"slide {index}: {shape.name}",
                    )
                    self.assertLessEqual(
                        shape.left + shape.width,
                        deck.slide_width - 457200,
                        f"slide {index}: {shape.name} violates the 0.5-inch right safe area",
                    )
                    self.assertLessEqual(
                        shape.top + shape.height,
                        deck.slide_height,
                        f"slide {index}: {shape.name}",
                    )
                editable_shapes = [
                    shape
                    for shape in slide.shapes
                    if shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX)
                    and getattr(shape, "has_text_frame", False)
                ]
                self.assertGreaterEqual(len(editable_shapes), 3, f"slide {index}")

            proposition_slide = slide_texts[9]
            for proposition in (f"P{number}" for number in range(1, 8)):
                self.assertIn(proposition, proposition_slide)
            self.assertIn("검증 전 명제", proposition_slide)

    def test_v03_pdf_export_has_twelve_pages_when_libreoffice_is_available(self):
        module_name = "tools.build_research_model_v03_media"
        self.assertIsNotNone(importlib.util.find_spec(module_name))
        build_media = importlib.import_module(module_name).build_media
        soffice = Path(r"C:\Program Files\LibreOffice\program\soffice.exe")
        if not soffice.is_file():
            self.skipTest("LibreOffice is not installed")

        with TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            deck_path = tmp_path / "deck.pptx"
            build_media(tmp_path / "model.png", deck_path)
            profile = tmp_path / "lo-profile"
            result = subprocess.run(
                [
                    str(soffice),
                    "--headless",
                    f"-env:UserInstallation={profile.resolve().as_uri()}",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(tmp_path),
                    str(deck_path),
                ],
                check=False,
                capture_output=True,
                text=True,
                timeout=120,
            )
            self.assertEqual(result.returncode, 0, result.stderr or result.stdout)
            pdf_path = deck_path.with_suffix(".pdf")
            self.assertTrue(pdf_path.is_file(), result.stderr or result.stdout)
            pages = PdfReader(pdf_path).pages
            self.assertEqual(len(pages), 12)
            for index, (page, contract) in enumerate(
                zip(pages, SLIDE_SEMANTIC_CONTRACT, strict=True),
                start=1,
            ):
                page_text = normalized_pdf_text(page)
                self.assertIn(contract["title"], page_text, f"PDF page {index} title")
                for fact in contract["facts"]:
                    self.assertIn(fact, page_text, f"PDF page {index} required fact")


if __name__ == "__main__":
    unittest.main()
