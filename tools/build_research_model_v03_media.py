"""Build the public v0.3 model diagram and editable advisor deck."""

from __future__ import annotations

import argparse
import os
import shutil
import struct
import subprocess
import zipfile
from datetime import datetime
from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


FONT = "Malgun Gothic"
INK = "111111"
MUTED = "5C6470"
PANEL = "EDEDED"
RULE = "B8BCC4"
BLUE = "3D8DFF"
BLUE_LIGHT = "DDEBFF"
AMBER = "F2B84B"
AMBER_LIGHT = "FFF0C9"
GREEN = "2E9F6B"
GREEN_LIGHT = "DDF3E8"
RED = "D95252"
RED_LIGHT = "F9E1E1"
TEAL = "188F96"
TEAL_LIGHT = "D9EFF0"
VIOLET = "7B61C8"
VIOLET_LIGHT = "E9E2F8"
WHITE = "FFFFFF"
SLIDE_WIDTH = Inches(13.333333)
SLIDE_HEIGHT = Inches(7.5)

ROOT = Path(__file__).resolve().parents[1]

SOURCE_URLS = {
    "kemell": "https://doi.org/10.1016/j.infsof.2025.107805",
    "neumann": "https://doi.org/10.1007/978-3-032-22375-3_18",
    "golgeci": "https://doi.org/10.1016/j.hrmr.2024.101075",
    "bc2012": "https://doi.org/10.5465/amj.2009.0891",
    "bc2013": "https://doi.org/10.1287/mnsc.1120.1583",
}


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _format_text_frame(
    text_frame,
    *,
    size: float,
    color: str = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margins: float = 0.08,
    wrap: bool = True,
) -> None:
    text_frame.word_wrap = wrap
    text_frame.vertical_anchor = valign
    text_frame.margin_left = Inches(margins)
    text_frame.margin_right = Inches(margins)
    text_frame.margin_top = Inches(margins)
    text_frame.margin_bottom = Inches(margins)
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.05
        for run in paragraph.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = _rgb(color)


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margins: float = 0.03,
    wrap: bool = True,
    name: str | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    if name:
        shape.name = name
    shape.text = text
    _format_text_frame(
        shape.text_frame,
        size=size,
        color=color,
        bold=bold,
        align=align,
        valign=valign,
        margins=margins,
        wrap=wrap,
    )
    return shape


def _add_box(
    slide,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    fill: str = PANEL,
    line: str = RULE,
    line_width: float = 1,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margins: float = 0.14,
    shape_type: MSO_SHAPE = MSO_SHAPE.RECTANGLE,
    name: str | None = None,
):
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    if name:
        shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line)
    shape.line.width = Pt(line_width)
    shape.text = text
    _format_text_frame(
        shape.text_frame,
        size=size,
        color=color,
        bold=bold,
        align=align,
        valign=valign,
        margins=margins,
    )
    return shape


def _add_rule(slide, x1: float, y1: float, x2: float, y2: float, *, color: str = RULE, width: float = 1):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = _rgb(color)
    line.line.width = Pt(width)
    return line


def _add_arrow(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = INK,
    width: float = 1.5,
    dashed: bool = False,
):
    connector = _add_rule(slide, x1, y1, x2, y2, color=color, width=width)
    if dashed:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    line_xml = connector.line._get_or_add_ln()
    tail_end = OxmlElement("a:tailEnd")
    tail_end.set("type", "triangle")
    tail_end.set("w", "sm")
    tail_end.set("len", "sm")
    line_xml.append(tail_end)
    return connector


def _blank_slide(presentation: Presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(WHITE)
    return slide


def _add_slide_shell(slide, title: str, number: int, section: str) -> None:
    _add_text(slide, section, 0.55, 0.16, 3.4, 0.22, size=16, color=MUTED, bold=True, wrap=False)
    _add_text(
        slide,
        title,
        0.55,
        0.42,
        12.2,
        0.55,
        size=36,
        bold=True,
        wrap=False,
        name="Slide Title",
    )
    _add_rule(slide, 0.55, 1.08, 12.78, 1.08, color=INK, width=1.2)
    _add_text(slide, "연구모형 v0.3 · 현장조사 전", 0.55, 7.12, 3.2, 0.18, size=10, color=MUTED, wrap=False)
    _add_text(slide, f"{number:02d}", 12.25, 7.12, 0.5, 0.18, size=10, color=MUTED, align=PP_ALIGN.RIGHT, wrap=False)


def _add_evidence_legend(slide, x: float = 6.1, y: float = 1.16) -> None:
    labels = (
        ("논문 직접 근거", BLUE),
        ("통합 해석", TEAL),
        ("검증 전 명제", VIOLET),
    )
    cursor = x
    for label, color in labels:
        _add_box(slide, "", cursor, y + 0.1, 0.12, 0.12, fill=color, line=color, line_width=0)
        _add_text(slide, label, cursor + 0.18, y, 1.72, 0.32, size=16, color=MUTED, wrap=False)
        cursor += 2.08


def _set_notes(slide, sources: list[str]) -> None:
    slide.notes_slide.notes_text_frame.text = "[Sources]\n" + "\n".join(
        f"- {source}" for source in sources
    )


def _diagram_dot() -> str:
    return f'''digraph model {{
  graph [rankdir=LR, bgcolor="#{WHITE}", margin=0, pad=0,
         size="10,6.666667!", ratio="fill", dpi=240,
         viewport="720,480,0.60",
         nodesep=0.42, ranksep=0.72, splines=ortho,
         fontname="{FONT}", labelloc="t",
         label="연구모형 v0.3 · 모든 경로는 검증 전 명제"];
  node [shape=box, style="filled", penwidth=2.2,
        fontname="{FONT}", fontsize=18, margin="0.18,0.12"];
  edge [color="#6B7280", penwidth=2.0, arrowsize=0.8,
        fontname="{FONT}", fontsize=13];

  role [label="실행된 역할 보완성\\n공식 권한 × 현장 기능", fillcolor="#{BLUE_LIGHT}", color="#{BLUE}"];
  network [label="활성화된 네트워크 연결성\\n조언 · 신뢰 · 승인/예외 · 위험 에스컬레이션", fillcolor="#{BLUE_LIGHT}", color="#{BLUE}"];

  safety [label="발언 안전성\\n말해도 불이익이 없는가?", fillcolor="#{AMBER_LIGHT}", color="#{AMBER}"];
  efficacy [label="발언 효능감\\n말하면 바뀔 수 있는가?", fillcolor="#{AMBER_LIGHT}", color="#{AMBER}"];

  surface [label="건설적 표면화\\n개선 발언 · 조건부 수용", fillcolor="#{GREEN_LIGHT}", color="#{GREEN}"];
  conceal [label="은폐 · 이탈\\n침묵 · 비사용 · 우회", fillcolor="#{RED_LIGHT}", color="#{RED}"];

  alignment [label="정책-실무 정합성", fillcolor="#{TEAL_LIGHT}", color="#{TEAL}"];
  adoption [label="책임 있는 AI 채택", fillcolor="#{TEAL_LIGHT}", color="#{TEAL}"];

  divergence [label="변화 발산성\\n업무 · 전문성 · 권한 · 책임", fillcolor="#{VIOLET_LIGHT}", color="#{VIOLET}"];

  {{rank=same; role; network}}
  {{rank=same; safety; efficacy}}
  {{rank=same; surface; conceal}}
  {{rank=same; alignment; adoption}}

  role -> safety;
  role -> efficacy;
  network -> safety;
  network -> efficacy;
  safety -> surface;
  efficacy -> surface;
  safety -> conceal;
  efficacy -> conceal;
  surface -> alignment [color="#{GREEN}"];
  surface -> adoption [color="#{GREEN}"];
  conceal -> alignment [color="#{RED}", style=dashed];
  conceal -> adoption [color="#{RED}", style=dashed];
  divergence -> network [color="#{VIOLET}", style=dashed, constraint=false];
  divergence -> surface [color="#{VIOLET}", style=dashed, constraint=false];
  divergence -> conceal [color="#{VIOLET}", style=dashed, constraint=false];
}}'''


def _find_dot() -> Path:
    discovered = shutil.which("dot")
    candidates = [
        Path(discovered) if discovered else None,
        Path(r"C:\Program Files\Graphviz\bin\dot.exe"),
        Path(r"C:\Program Files (x86)\Graphviz\bin\dot.exe"),
    ]
    for candidate in candidates:
        if candidate and candidate.is_file():
            return candidate
    raise RuntimeError("Graphviz dot was not found. Install Graphviz.Graphviz before building media.")


def _read_png_size(path: Path) -> tuple[int, int]:
    with path.open("rb") as stream:
        header = stream.read(24)
    if header[:8] != b"\x89PNG\r\n\x1a\n":
        raise RuntimeError(f"Graphviz did not produce a PNG: {path}")
    return struct.unpack(">II", header[16:24])


def _build_diagram(diagram_path: Path) -> None:
    diagram_path.parent.mkdir(parents=True, exist_ok=True)
    with TemporaryDirectory() as tmp:
        dot_path = Path(tmp) / "research-model-v0.3.dot"
        dot_path.write_text(_diagram_dot(), encoding="utf-8")
        result = subprocess.run(
            [str(_find_dot()), "-Tpng:cairo", str(dot_path), "-o", str(diagram_path)],
            check=False,
            capture_output=True,
            text=True,
            encoding="utf-8",
        )
    if result.returncode:
        raise RuntimeError(result.stderr.strip() or "Graphviz diagram build failed")
    actual_size = _read_png_size(diagram_path)
    if actual_size != (2400, 1600):
        raise RuntimeError(f"Expected a 2400x1600 Graphviz PNG, got {actual_size}")


def _slide_1(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_text(slide, "지도교수 세미나 · 연구모형 v0.3", 0.65, 0.55, 5.3, 0.35, size=18, color=MUTED, bold=True, wrap=False)
    _add_text(
        slide,
        "분산된 AI 변화주도자 네트워크",
        0.65,
        1.25,
        12.0,
        0.85,
        size=54,
        bold=True,
        wrap=False,
        name="Slide Title",
    )
    _add_rule(slide, 0.65, 2.35, 12.65, 2.35, color=INK, width=2)
    _add_text(
        slide,
        "공식 리더와 비공식 AI 챔피언의 역할 보완성과 네트워크 연결성이\n"
        "AI 사건의 발언 판단과 반응 경로를 어떻게 바꾸는지 탐색한다.",
        0.65,
        2.75,
        10.8,
        1.25,
        size=26,
        color=INK,
        valign=MSO_ANCHOR.TOP,
    )
    _add_box(
        slide,
        "현장조사 전 · P1-P7은 검증 전 명제",
        0.65,
        5.55,
        5.35,
        0.58,
        fill=VIOLET_LIGHT,
        line=VIOLET,
        size=18,
        color=VIOLET,
        bold=True,
    )
    _add_text(slide, "2026.08.12", 10.85, 6.75, 1.8, 0.3, size=16, color=MUTED, align=PP_ALIGN.RIGHT, wrap=False)
    _set_notes(slide, ["site/downloads/research-design/research-model-v0.3.md"])


def _slide_2(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_slide_shell(slide, "공식 도입과 비공식 실험이 같은 사건에 겹친다", 2, "01 · 현상과 연구 공백")
    _add_arrow(slide, 2.55, 2.33, 4.55, 3.05, color=BLUE)
    _add_arrow(slide, 2.55, 4.35, 4.55, 3.55, color=BLUE)
    _add_box(slide, "공식 도입\n승인 도구 · 정책 · 책임", 0.75, 1.75, 2.5, 1.15, fill=BLUE_LIGHT, line=BLUE, size=20, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "비공식 실험\n상향식 수요 · 우회 · Shadow AI", 0.75, 3.85, 2.5, 1.15, fill=PANEL, line=RULE, size=20, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "같은 AI 사건", 4.55, 2.55, 2.2, 1.05, fill=WHITE, line=INK, line_width=2, size=24, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "논문 직접 근거", 7.35, 1.55, 2.1, 0.35, size=18, color=BLUE, bold=True, wrap=False)
    _add_text(slide, "SW 기업에서 공식 제공과 현장 실험이 병존하고,\n정책-실무 간극과 비승인 사용 신호가 관찰됐다.", 7.35, 1.95, 5.05, 1.05, size=19, valign=MSO_ANCHOR.TOP)
    _add_text(slide, "연구 공백", 7.35, 3.3, 2.1, 0.35, size=18, color=TEAL, bold=True, wrap=False)
    _add_text(slide, "누가 어떤 역할을 맡고, 어떤 관계를 통해 우려가\n도달하며, 왜 표면화 또는 은폐되는지는 설명이 부족하다.", 7.35, 3.7, 5.05, 1.15, size=21, bold=True, valign=MSO_ANCHOR.TOP)
    _add_box(slide, "통합 해석: 분석의 초점은 일반 태도가 아니라 구체적 AI 사건이다.", 7.35, 5.35, 5.05, 0.7, fill=TEAL_LIGHT, line=TEAL, size=17, color=INK)
    _set_notes(slide, [
        "site/downloads/kemell-2025/model-v0.3-contribution.md",
        SOURCE_URLS["kemell"],
        "site/downloads/neumann-2026/model-v0.3-contribution.md",
        SOURCE_URLS["neumann"],
        "site/downloads/research-design/research-model-v0.3.md",
    ])


def _slide_3(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_slide_shell(slide, "다섯 논문은 현상·과정·경계의 서로 다른 근거다", 3, "01 · 증거 아키텍처")
    _add_evidence_legend(slide)
    _add_text(slide, "논문", 0.7, 1.65, 2.55, 0.34, size=17, color=MUTED, bold=True)
    _add_text(slide, "직접 기여 범위", 3.25, 1.65, 4.1, 0.34, size=17, color=MUTED, bold=True)
    _add_text(slide, "v0.3에서의 연결", 7.65, 1.65, 4.75, 0.34, size=17, color=MUTED, bold=True)
    rows = [
        ("Kemell et al. (2025)", "공식 도입과 현장 실험의 병존", "현상 앵커 · P1-P5 표집 근거"),
        ("Neumann et al. (2026)", "정책-실무 간극과 Shadow AI 신호", "은폐·이탈 및 정합성 맥락"),
        ("Golgeci et al. (2025)", "AI 저항·완화의 개념적 과정", "발언 판단과 반응 경로 해석"),
        ("Battilana & Casciaro (2012)", "발산성별 네트워크 조건", "P6의 조건부 경계"),
        ("Battilana & Casciaro (2013)", "수신자 태도·관계강도 경계", "P7의 설득 경계"),
    ]
    for index, (paper, direct, synthesis) in enumerate(rows):
        y = 2.05 + index * 0.82
        if index % 2 == 0:
            _add_box(slide, "", 0.6, y - 0.05, 12.15, 0.72, fill="F6F6F6", line="F6F6F6", line_width=0)
        _add_text(slide, paper, 0.72, y, 2.55, 0.58, size=16, bold=True)
        _add_text(slide, direct, 3.25, y, 4.1, 0.58, size=16)
        _add_text(slide, synthesis, 7.65, y, 4.65, 0.58, size=16, color=TEAL)
    _add_text(slide, "경계", 0.72, 6.35, 0.75, 0.32, size=17, color=RED, bold=True, wrap=False)
    _add_text(slide, "다섯 논문을 합쳐도 P1-P7의 직접 검증 근거가 되지 않는다.", 1.55, 6.28, 9.8, 0.42, size=18, bold=True, wrap=False)
    _set_notes(slide, [
        "site/downloads/kemell-2025/model-v0.3-contribution.md", SOURCE_URLS["kemell"],
        "site/downloads/neumann-2026/model-v0.3-contribution.md", SOURCE_URLS["neumann"],
        "site/downloads/golgeci-2025/model-v0.3-contribution.md", SOURCE_URLS["golgeci"],
        "site/downloads/battilana-casciaro-2012/model-v0.3-contribution.md", SOURCE_URLS["bc2012"],
        "site/downloads/battilana-casciaro-2013/model-v0.3-contribution.md", SOURCE_URLS["bc2013"],
    ])


def _slide_4(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_slide_shell(slide, "AI 사건을 중심으로 세 분석 수준을 분리한다", 4, "02 · 초점 사건과 분석 단위")
    _add_arrow(slide, 2.5, 3.62, 3.25, 3.62, color=INK)
    _add_arrow(slide, 5.0, 3.62, 5.75, 3.62, color=INK)
    _add_arrow(slide, 7.5, 3.62, 8.25, 3.62, color=INK)
    _add_arrow(slide, 10.0, 3.62, 10.75, 3.62, color=INK)
    _add_box(slide, "팀 맥락", 0.72, 1.45, 2.0, 0.55, fill=PANEL, line=RULE, size=20, bold=True)
    _add_text(slide, "기초 네트워크 구조\n정적 중심성 · 브로커리지 · 응집성", 0.72, 2.05, 3.15, 0.95, size=18, valign=MSO_ANCHOR.TOP)
    _add_box(slide, "사건", 4.58, 1.45, 1.4, 0.55, fill=BLUE_LIGHT, line=BLUE, size=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "역할 보완성 · 연결성 · 변화 발산성\n지배 경로 · 후행 과정결과", 4.05, 2.05, 4.1, 0.95, size=19, bold=True, valign=MSO_ANCHOR.TOP, align=PP_ALIGN.CENTER)
    _add_box(slide, "개인 × 사건", 9.82, 1.45, 2.0, 0.55, fill=AMBER_LIGHT, line=AMBER, size=20, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "발언 안전성 · 발언 효능감\n복수 구성원 반응", 9.35, 2.05, 3.25, 0.95, size=18, valign=MSO_ANCHOR.TOP, align=PP_ALIGN.CENTER)
    phases = ("촉발", "관계 동원", "발언 판단", "복수 반응", "후행 결과")
    for index, phase in enumerate(phases):
        _add_box(slide, phase, 0.75 + index * 2.5, 3.22, 1.75, 0.8, fill=WHITE, line=INK if index else BLUE, line_width=2 if index == 0 else 1.2, size=18, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "단위 계약: 전달 성공은 연결성, 실제 문제해결은 후행 결과로 분리한다.", 2.3, 5.2, 8.7, 0.72, fill=TEAL_LIGHT, line=TEAL, size=19, bold=True, align=PP_ALIGN.CENTER)
    _set_notes(slide, [
        "site/downloads/research-design/research-model-v0.3.md",
        "site/downloads/research-design/construct-dictionary-v0.3.md",
        "site/downloads/research-design/pilot-protocol-and-codingbook-v0.3.md",
    ])


def _slide_5(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_slide_shell(slide, "역할 보완성은 권한과 현장 기능의 실제 분담이다", 5, "02 · 핵심 예측변수")
    _add_arrow(slide, 3.95, 3.65, 5.15, 3.65, color=BLUE, width=2)
    _add_arrow(slide, 9.35, 3.65, 8.15, 3.65, color=BLUE, width=2)
    _add_text(slide, "공식 리더", 0.75, 1.55, 3.2, 0.42, size=24, bold=True, color=BLUE)
    _add_text(slide, "방향·허용범위\n자원·승인·예외\n최종 책임", 0.75, 2.15, 3.2, 2.3, size=21, valign=MSO_ANCHOR.TOP)
    _add_box(slide, "실행된\n역할 보완성", 5.15, 2.65, 3.0, 2.0, fill=BLUE_LIGHT, line=BLUE, line_width=2.5, size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "비공식 AI 챔피언", 9.35, 1.55, 3.2, 0.42, size=24, bold=True, color=BLUE)
    _add_text(slide, "현장 실험·번역\n동료 코칭\n위험 발견·전달", 9.35, 2.15, 3.2, 2.3, size=21, valign=MSO_ANCHOR.TOP)
    _add_rule(slide, 0.75, 5.35, 12.55, 5.35, color=RULE, width=1.2)
    _add_text(slide, "포함", 0.75, 5.62, 0.75, 0.3, size=18, color=GREEN, bold=True, wrap=False)
    _add_text(slide, "사건에서 실제 실행된 기능·책임·에스컬레이션 경로", 1.55, 5.55, 5.1, 0.5, size=18)
    _add_text(slide, "제외", 7.0, 5.62, 0.75, 0.3, size=18, color=RED, bold=True, wrap=False)
    _add_text(slide, "직급·친분·개인 열의 또는 해결 성과", 7.8, 5.55, 4.5, 0.5, size=18)
    _set_notes(slide, [
        "site/downloads/research-design/construct-dictionary-v0.3.md",
        "site/downloads/research-design/research-model-v0.3.md",
        "site/downloads/golgeci-2025/model-v0.3-contribution.md",
        SOURCE_URLS["golgeci"],
    ])


def _slide_6(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_slide_shell(slide, "네트워크 연결성은 네 관계의 도달·확인을 묻는다", 6, "02 · 다중 관계망")
    for index in range(4):
        y = 1.82 + index * 1.08
        _add_arrow(slide, 2.6, y + 0.34, 4.0, y + 0.34, color=BLUE, width=1.8)
    relations = (
        ("조언", "문제 해결 지식과 사용 방법"),
        ("신뢰", "민감한 우려를 맡길 관계"),
        ("승인·예외", "허용범위와 책임을 결정할 경로"),
        ("위험 에스컬레이션", "보안·품질·법무 위험의 상향 전달"),
    )
    for index, (relation, meaning) in enumerate(relations):
        y = 1.82 + index * 1.08
        _add_box(slide, relation, 0.72, y, 1.88, 0.68, fill=BLUE_LIGHT if index < 2 else WHITE, line=BLUE, size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, meaning, 4.0, y, 3.85, 0.68, size=18)
    _add_text(slide, "관계별 증거 사슬", 8.25, 1.65, 3.2, 0.38, size=24, bold=True)
    _add_arrow(slide, 9.4, 2.8, 9.4, 3.25, color=INK)
    _add_arrow(slide, 9.4, 4.0, 9.4, 4.45, color=INK)
    _add_box(slide, "발신자 보고 전달", 8.15, 2.15, 2.5, 0.65, fill=PANEL, line=RULE, size=18, align=PP_ALIGN.CENTER)
    _add_box(slide, "수신·확인 근거", 8.15, 3.25, 2.5, 0.75, fill=BLUE_LIGHT, line=BLUE, size=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "수신 뒤 태도", 8.15, 4.45, 2.5, 0.65, fill=PANEL, line=RULE, size=18, align=PP_ALIGN.CENTER)
    _add_text(slide, "발신자 보고만으로 수신을 확정하지 않는다.", 7.75, 5.55, 4.3, 0.55, size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)
    _set_notes(slide, [
        "site/downloads/research-design/construct-dictionary-v0.3.md",
        "site/downloads/research-design/pilot-protocol-and-codingbook-v0.3.md",
        "tools/build_research_model_v03_workbook.py",
    ])


def _slide_7(presentation: Presentation, diagram_path: Path) -> None:
    slide = _blank_slide(presentation)
    _add_slide_shell(slide, "안전성과 효능감은 서로 다른 발언 판단이다", 7, "03 · 이중 매개 판단")
    _add_box(slide, "발언 안전성", 0.72, 1.65, 3.2, 0.72, fill=AMBER_LIGHT, line=AMBER, size=24, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "“이 사건을 말해도\n불이익이 없는가?”", 0.9, 2.55, 2.85, 1.05, size=22, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "발언 효능감", 0.72, 4.15, 3.2, 0.72, fill=AMBER_LIGHT, line=AMBER, size=24, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "“말하면 권한 있는 행위자가\n무언가를 바꿀 수 있는가?”", 0.88, 5.05, 2.9, 1.05, size=20, bold=True, align=PP_ALIGN.CENTER)
    picture = slide.shapes.add_picture(
        str(diagram_path), Inches(4.35), Inches(1.42), width=Inches(8.1), height=Inches(5.4)
    )
    picture.name = "Research model diagram"
    picture._element.nvPicPr.cNvPr.set("descr", "research-model-v0.3.png")
    _add_text(slide, "둘은 함께 움직일 수 있지만 같은 구성개념이 아니다.", 4.45, 6.42, 7.9, 0.35, size=18, color=AMBER, bold=True, align=PP_ALIGN.CENTER, wrap=False)
    _set_notes(slide, [
        "site/downloads/research-design/research-model-v0.3.png",
        "site/downloads/research-design/research-model-v0.3.md",
        "site/downloads/research-design/construct-dictionary-v0.3.md",
    ])


def _slide_8(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_slide_shell(slide, "한 사건의 반응은 복수 행동과 지배 경로로 코딩한다", 8, "03 · 반응 경로")
    _add_arrow(slide, 2.5, 2.25, 3.0, 2.25, color=INK)
    _add_arrow(slide, 5.55, 2.25, 6.05, 2.25, color=INK)
    _add_arrow(slide, 8.6, 2.25, 9.1, 2.25, color=INK)
    _add_box(slide, "촉발·위협", 0.75, 1.82, 1.75, 0.85, fill=PANEL, line=RULE, size=20, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "개인 반응", 3.0, 1.82, 2.55, 0.85, fill=AMBER_LIGHT, line=AMBER, size=20, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "사건 지배 경로", 6.05, 1.82, 2.55, 0.85, fill=WHITE, line=INK, size=20, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "후행 과정결과", 9.1, 1.82, 2.75, 0.85, fill=TEAL_LIGHT, line=TEAL, size=20, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "개선 발언 · 조건부 수용\n침묵 · 비사용 · 예외 요청\n우회 · 비승인 사용 · 공개 반대", 3.0, 2.95, 2.65, 1.65, size=17, valign=MSO_ANCHOR.TOP)
    _add_box(slide, "건설적 표면화", 6.05, 3.0, 2.55, 0.7, fill=GREEN_LIGHT, line=GREEN, size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "은폐·이탈", 6.05, 4.0, 2.55, 0.7, fill=RED_LIGHT, line=RED, size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "검토 착수·조정과\n실질 처리·문제해결을 분리", 9.1, 3.1, 2.75, 1.2, size=19, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "우회와 Shadow AI는 의도·필요성·위험·공식 대안을 확인한 뒤 판정한다.", 1.2, 5.55, 10.9, 0.72, fill=PANEL, line=RULE, size=18, bold=True, align=PP_ALIGN.CENTER)
    _set_notes(slide, [
        "site/downloads/research-design/research-model-v0.3.md",
        "site/downloads/research-design/construct-dictionary-v0.3.md",
        "site/downloads/research-design/pilot-protocol-and-codingbook-v0.3.md",
        "tools/build_research_model_v03_workbook.py",
    ])


def _slide_9(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_slide_shell(slide, "변화 발산성이 유리한 관계 구조를 바꿀 수 있다", 9, "04 · 경계조건")
    _add_evidence_legend(slide)
    _add_arrow(slide, 1.15, 5.2, 11.95, 5.2, color=VIOLET, width=2.5)
    _add_text(slide, "낮은 변화 발산성", 0.78, 5.45, 3.2, 0.4, size=20, color=VIOLET, bold=True, wrap=False)
    _add_text(slide, "높은 변화 발산성", 9.35, 5.45, 3.0, 0.4, size=20, color=VIOLET, bold=True, align=PP_ALIGN.RIGHT, wrap=False)
    _add_text(slide, "기존 관행에 가까움", 0.78, 5.9, 2.8, 0.35, size=17, color=MUTED, wrap=False)
    _add_text(slide, "업무·전문성·권한·책임 재배치", 8.65, 5.9, 3.7, 0.35, size=17, color=MUTED, align=PP_ALIGN.RIGHT, wrap=False)
    _add_box(slide, "내부 응집·강한 관계", 1.0, 2.0, 4.35, 1.25, fill=PANEL, line=RULE, size=25, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "경계 연결·중개", 7.95, 2.0, 4.35, 1.25, fill=VIOLET_LIGHT, line=VIOLET, size=25, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "논문 직접 근거: NHS 변화 프로젝트에서\n발산성에 따른 조건부 네트워크 관계가 보고됐다.", 0.98, 3.55, 4.45, 0.9, size=17, valign=MSO_ANCHOR.TOP)
    _add_text(slide, "검증 전 명제: SW·AI 다중 행위자 맥락으로\n확장한 P6은 연구 1에서 수정·분기·기각한다.", 7.95, 3.55, 4.4, 0.9, size=17, valign=MSO_ANCHOR.TOP)
    _add_box(slide, "P7 경계: 고발산 사건의 반복 설득은 관계압력·손상·소진 가능성도 함께 본다.", 2.0, 6.35, 9.35, 0.55, fill=RED_LIGHT, line=RED, size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)
    _set_notes(slide, [
        "site/downloads/battilana-casciaro-2012/model-v0.3-contribution.md", SOURCE_URLS["bc2012"],
        "site/downloads/battilana-casciaro-2013/model-v0.3-contribution.md", SOURCE_URLS["bc2013"],
        "site/downloads/research-design/proposition-traceability-v0.3.md",
    ])


def _slide_10(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_slide_shell(slide, "P1-P7은 경쟁 설명과 함께 검증할 명제다", 10, "05 · 명제와 반증")
    _add_evidence_legend(slide)
    _add_text(slide, "명제", 0.7, 1.62, 0.75, 0.3, size=17, color=MUTED, bold=True)
    _add_text(slide, "탐색할 관계", 1.6, 1.62, 6.25, 0.3, size=17, color=MUTED, bold=True)
    _add_text(slide, "대표 경쟁 설명", 8.15, 1.62, 4.25, 0.3, size=17, color=MUTED, bold=True)
    rows = (
        ("P1", "보완성·연결성 → 안전성·효능감", "기존 리더 신뢰 · AI 경험 · 역인과"),
        ("P2", "발언 안전성 → 우려의 표면화", "발언 의무 · 외부 감사"),
        ("P3", "발언 효능감 → 개선 요구", "개인 전문성 · 이미 확보된 자원"),
        ("P4", "낮은 판단 → 은폐·이탈", "대안 부재 · 긴급 과업 · 도구 접근성"),
        ("P5", "표면화·처리 → 원거리 결과", "법무·고객 요구 · 예산 · 기술 성숙"),
        ("P6", "발산성별 중개·응집 조건", "과업 위험 · 권한 · 조직 규모"),
        ("P7", "반복 설득의 관계 경계", "기존 갈등 · 과부하 · 권력차"),
    )
    for index, (proposition, relation, rival) in enumerate(rows):
        y = 2.0 + index * 0.62
        if index % 2 == 0:
            _add_box(slide, "", 0.6, y - 0.02, 12.15, 0.56, fill="F6F6F6", line="F6F6F6", line_width=0)
        _add_text(slide, proposition, 0.72, y, 0.65, 0.5, size=17, color=VIOLET, bold=True)
        _add_text(slide, relation, 1.6, y, 6.2, 0.5, size=16)
        _add_text(slide, rival, 8.15, y, 4.25, 0.5, size=16, color=MUTED)
    _add_box(slide, "판정 규칙: 지지 사례만 찾지 않고 반증 패턴을 같은 표에 기록한다.", 2.0, 6.48, 9.3, 0.48, fill=VIOLET_LIGHT, line=VIOLET, size=16, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
    _set_notes(slide, [
        "site/downloads/research-design/proposition-traceability-v0.3.md",
        "site/downloads/research-design/research-model-v0.3.md",
        "site/downloads/research-design/construct-dictionary-v0.3.md",
    ])


def _slide_11(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_slide_shell(slide, "탐색적 순차 혼합연구가 명제를 단계적으로 거른다", 11, "06 · 검증 설계")
    _add_arrow(slide, 4.05, 3.5, 4.55, 3.5, color=INK, width=2)
    _add_arrow(slide, 8.35, 3.5, 8.85, 3.5, color=INK, width=2)
    _add_text(slide, "연구 1", 0.78, 1.55, 1.3, 0.35, size=20, color=BLUE, bold=True, wrap=False)
    _add_text(slide, "연구 2", 4.58, 1.55, 1.3, 0.35, size=20, color=AMBER, bold=True, wrap=False)
    _add_text(slide, "연구 3", 8.88, 1.55, 1.3, 0.35, size=20, color=TEAL, bold=True, wrap=False)
    _add_box(slide, "프레임워크 개발", 0.75, 2.0, 3.3, 1.0, fill=BLUE_LIGHT, line=BLUE, size=25, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "측정모형 정제", 4.55, 2.0, 3.8, 1.0, fill=AMBER_LIGHT, line=AMBER, size=25, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "가설 검증", 8.85, 2.0, 3.7, 1.0, fill=TEAL_LIGHT, line=TEAL, size=25, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "결정적 사건 재구성\n시간 순서 · 반례 · 경쟁 설명\n개념 경계 수정", 0.85, 3.35, 3.05, 1.65, size=18, valign=MSO_ANCHOR.TOP)
    _add_text(slide, "문항 풀 · 반응경로 코딩\n내용타당도 · 인지면접 · 차원성\nreflective/formative/index 결정", 4.67, 3.35, 3.55, 1.65, size=17, valign=MSO_ANCHOR.TOP)
    _add_text(slide, "유지된 경로만 H1-Hn 전환\n다원자료 · 시간적 선후 · 역인과\n검정력 분석 뒤 표본·기법 결정", 8.98, 3.35, 3.42, 1.65, size=17, valign=MSO_ANCHOR.TOP)
    _add_box(slide, "Gate 1: 개념·시간순서·반례", 1.15, 5.55, 3.15, 0.55, fill=WHITE, line=BLUE, size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "Gate 2: 측정 타당성", 5.0, 5.55, 3.0, 0.55, fill=WHITE, line=AMBER, size=16, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "사전등록 후 검증", 9.25, 5.55, 2.85, 0.55, fill=WHITE, line=TEAL, size=16, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "현 단계에서는 특정 표본 수나 SEM 유형을 고정하지 않는다.", 2.6, 6.45, 8.2, 0.38, size=18, color=RED, bold=True, align=PP_ALIGN.CENTER, wrap=False)
    _set_notes(slide, [
        "site/downloads/research-design/research-model-v0.3.md",
        "site/downloads/research-design/pilot-protocol-and-codingbook-v0.3.md",
        "site/downloads/research-design/construct-dictionary-v0.3.md",
        "site/downloads/research-design/proposition-traceability-v0.3.md",
    ])


def _slide_12(presentation: Presentation) -> None:
    slide = _blank_slide(presentation)
    _add_slide_shell(slide, "현장 진입 전 윤리 게이트와 다음 결정을 확인한다", 12, "07 · 결정")
    _add_text(slide, "현장 진입 전 필수 게이트", 0.75, 1.45, 5.55, 0.45, size=24, bold=True)
    gates = (
        "IRB + 기업 보안·법무 승인",
        "독립 모집·동의 회수와 원자료 비공유",
        "연결코드 분리 암호화·접근·폐기",
        "작은 팀 재식별 방지와 비응답 지명자 보호",
        "내부자 위치성·이해상충 공개",
    )
    for index, gate in enumerate(gates):
        y = 2.05 + index * 0.72
        _add_text(slide, f"{index + 1:02d}", 0.78, y, 0.48, 0.42, size=17, color=BLUE, bold=True, wrap=False)
        _add_text(slide, gate, 1.4, y, 5.2, 0.42, size=17, bold=index == 0)
        _add_rule(slide, 1.38, y + 0.48, 6.35, y + 0.48, color=RULE)
    _add_text(slide, "현재 검토 상태", 7.15, 1.45, 4.8, 0.45, size=24, bold=True)
    _add_box(slide, "설계 게이트 PASS\nCritical 0 · High 0", 7.15, 2.05, 5.1, 1.05, fill=GREEN_LIGHT, line=GREEN, size=23, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    _add_box(slide, "현장조사는 시작하지 않았다\nNotebookLM 산출물은 근거가 아니다", 7.15, 3.38, 5.1, 1.05, fill=PANEL, line=RULE, size=19, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "다음 결정", 7.15, 4.88, 1.6, 0.35, size=20, color=VIOLET, bold=True, wrap=False)
    _add_text(slide, "IRB·기업 승인·모집·자료접근 문서 준비와\n60분 파일럿 설계 검토를 진행할 것인가?", 7.15, 5.32, 5.25, 1.05, size=22, bold=True, valign=MSO_ANCHOR.TOP)
    _set_notes(slide, [
        "docs/reviews/2026-08-12-research-model-v0.3-review.md",
        "docs/superpowers/specs/2026-08-12-research-model-v0.3-design.md",
        "site/downloads/research-design/pilot-protocol-and-codingbook-v0.3.md",
        "site/downloads/research-design/research-model-v0.3.md",
    ])


def _normalize_pptx(path: Path) -> None:
    normalized_path = path.with_suffix(".normalized.pptx")
    fixed_time = (2026, 8, 12, 0, 0, 0)
    with zipfile.ZipFile(path, "r") as source, zipfile.ZipFile(
        normalized_path, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9
    ) as target:
        for name in sorted(source.namelist()):
            source_info = source.getinfo(name)
            info = zipfile.ZipInfo(name, fixed_time)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = source_info.external_attr
            info.create_system = source_info.create_system
            target.writestr(info, source.read(name))
    os.replace(normalized_path, path)


def _build_deck(diagram_path: Path, deck_path: Path) -> None:
    deck_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    presentation.core_properties.title = "SW AI change-agent research model v0.3 - advisor deck"
    presentation.core_properties.subject = "Pre-fieldwork research model and validation design"
    presentation.core_properties.author = "Doctoral Paper Lab"
    presentation.core_properties.created = datetime(2026, 8, 12)
    presentation.core_properties.modified = datetime(2026, 8, 12)
    _slide_1(presentation)
    _slide_2(presentation)
    _slide_3(presentation)
    _slide_4(presentation)
    _slide_5(presentation)
    _slide_6(presentation)
    _slide_7(presentation, diagram_path)
    _slide_8(presentation)
    _slide_9(presentation)
    _slide_10(presentation)
    _slide_11(presentation)
    _slide_12(presentation)
    presentation.save(deck_path)
    _normalize_pptx(deck_path)


def build_media(diagram_path: Path, deck_path: Path) -> None:
    """Build one 2400x1600 Graphviz PNG and one editable 12-slide PPTX."""

    diagram_path = Path(diagram_path)
    deck_path = Path(deck_path)
    _build_diagram(diagram_path)
    _build_deck(diagram_path, deck_path)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--diagram", required=True, type=Path)
    parser.add_argument("--deck", required=True, type=Path)
    args = parser.parse_args()
    build_media(args.diagram, args.deck)


if __name__ == "__main__":
    main()
